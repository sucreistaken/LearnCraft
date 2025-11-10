import sys, json, os, traceback
from typing import List, Dict, Any

# --- Konuşma Tanıma (ASR) ---
try:
    from faster_whisper import WhisperModel
except Exception as e:
    WhisperModel = None  # import hatasını sonra yakalayıp anlaşılır mesaj veririz

# --- PDF / PPTX metin çıkarımı ---
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def transcribe(video_path: str) -> Dict[str, Any]:
    if WhisperModel is None:
        raise RuntimeError(
            "faster-whisper yüklü değil. `pip install faster-whisper` ile yükleyin."
        )

    model_size = os.getenv("WHISPER_MODEL", "tiny")  # tiny/small/medium/large-v3
    device = "cuda" if os.getenv("CUDA", "0") == "1" else "cpu"

    # not: CTranslate2 modeli ilk çağrıda indirir; internet gerekebilir
    model = WhisperModel(model_size, device=device)

    segments, info = model.transcribe(
        video_path,
        language="tr",
        beam_size=5,
        vad_filter=True
    )

    out_segments = []
    for s in segments:
        out_segments.append({
            "start": round(float(s.start), 2),
            "end": round(float(s.end), 2),
            "text": (s.text or "").strip()
        })

    return {
        "language": info.language,
        "duration": float(info.duration) if info.duration is not None else None,
        "segments": out_segments
    }


def extract_pdf(pdf_path: str) -> Dict[str, Any]:
    if fitz is None:
        raise RuntimeError("PyMuPDF (fitz) yüklü değil. `pip install pymupdf`")
    doc = fitz.open(pdf_path)
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text") or ""
        pages.append({"page": i + 1, "text": text.strip()})
    return {"type": "pdf", "pages": pages}


def extract_pptx(pptx_path: str) -> Dict[str, Any]:
    if Presentation is None:
        raise RuntimeError("python-pptx yüklü değil. `pip install python-pptx`")
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
        slides.append({"slide": i + 1, "text": "\n".join(texts)})
    return {"type": "pptx", "slides": slides}


def extract_slides(path_: str) -> Dict[str, Any]:
    low = path_.lower()
    if low.endswith(".pdf"):
        return extract_pdf(path_)
    if low.endswith(".pptx"):
        return extract_pptx(path_)
    return {"type": "unknown", "rawPath": path_}


def main() -> None:
    if len(sys.argv) < 3:
        print(json.dumps({"ok": False, "error": "usage: process_lecture.py <video> <slides>"}))
        sys.exit(2)

    video_path = sys.argv[1]
    slides_path = sys.argv[2]

    try:
        transcript = transcribe(video_path)
        slides = extract_slides(slides_path)
        result = {"ok": True, "transcript": transcript, "slides": slides}
        print(json.dumps(result, ensure_ascii=False))
    except Exception as e:
        err = {
            "ok": False,
            "error": str(e),
            "trace": traceback.format_exc().splitlines()[-5:]
        }
        print(json.dumps(err, ensure_ascii=False))
        sys.exit(1)


if __name__ == "__main__":
    main()
